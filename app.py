
from flask import Flask,request, send_file
from pdf2docx import Converter
import os
import tempfile
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from PIL import Image
import img2pdf
from werkzeug.utils import secure_filename


app = Flask(__name__)


# pdf转换为word
@app.route('/pdf_to_word', methods=['POST'])
def convert_pdf_to_word():
    # 假设前端通过表单上传了PDF文件
    pdf_file = request.files['pdf']  
    # 创建临时文件保存PDF文件
    pdf_temp = tempfile.NamedTemporaryFile(delete=False)
    pdf_file.save(pdf_temp.name)
    word_temp = None
    cv = None
    try:
        # 创建临时文件保存Word文档
        word_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.docx', mode='w+b')       
        # 使用pdf2docx库转换PDF为Word
        cv = Converter(pdf_temp.name)
        cv.convert(word_temp.name, start=0, end=None)    
        # 返回Word文档
        word_temp.seek(0)
        return send_file(word_temp.name, as_attachment=True, download_name='converted.docx')
    except Exception as e:
        print(f"转换失败: {e}")
        # 可以在这里添加更多的错误处理逻辑
    finally:
        # 删除临时的PDF和Word文件
        if cv:
            cv.close()
        if pdf_temp:
            os.unlink(pdf_temp.name)
        if word_temp:
            os.unlink(word_temp.name)

ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'bmp'}

def allowed_file(filename):
    """检查文件扩展名是否允许"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/image_to_pdf', methods=['POST'])
def convert_images_to_pdf():
    # 检查是否有文件被上传
    if 'images' not in request.files:
        return "没有文件部分", 400
    
    files = request.files.getlist('images')
    image_paths = []
    
    for file in files:
        if file and allowed_file(file.filename):
            # 保存图片到临时文件
            filename = secure_filename(file.filename)
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png', mode='w+b')
            file.save(temp_file.name)
            image_paths.append(temp_file.name)
        else:
            return "文件类型不允许", 400
    
    try:
        # 创建临时文件保存PDF文档
        pdf_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf', mode='w+b')    
        # 获取用户设置的页面方向和边距
        orientation = request.form.get('orientation', 'portrait')  # 默认为纵向
        margin = int(request.form.get('margin', 8))  # 默认边距为0    
        # 设置PDF的布局
        a4inpt = (img2pdf.mm_to_pt(210), img2pdf.mm_to_pt(297))  # A4大小，单位转换为点
        layout_fun = img2pdf.get_layout_fun(pagesize=a4inpt)       
        if orientation == 'landscape':
            # 如果用户选择横向，交换页面尺寸的宽高
            a4inpt = (a4inpt[1], a4inpt[0])  
        # 使用img2pdf库将所有图片合并为一个PDF
        with open(pdf_temp.name, "wb") as f:
            f.write(img2pdf.convert(image_paths, layout_fun=layout_fun))    
        # 返回PDF文档
        pdf_temp.seek(0)
        return send_file(pdf_temp.name, as_attachment=True, download_name='converted.pdf')
    finally:
        # 删除临时的图片文件和PDF文件
        for path in image_paths:
            os.unlink(path)
        os.unlink(pdf_temp.name)

@app.route('/pdf_to_pptx', methods=['POST'])
def convert_pdf_to_pptx():
    # 假设前端通过表单上传了PDF文件
    pdf_file = request.files['pdf']
    # 创建临时文件保存PDF
    pdf_temp = tempfile.NamedTemporaryFile(delete=False)
    pdf_file.save(pdf_temp.name)   
    pptx_temp = None  # 初始化pptx_temp变量
    try:
        # 创建PPT对象
        ppt = Presentation()   
        # 使用pdfplumber打开PDF文件
        with pdfplumber.open(pdf_temp.name) as pdf:
            # 遍历PDF中的每一页
            for page in pdf.pages:
                # 检查页面是否包含文本
                text = page.extract_text()
                if text:
                    # 为文本创建新的幻灯片
                    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                    tf = txBox.text_frame
                    p = tf.add_paragraph()
                    p.text = text
                    p.font.size = Pt(12)              
                # 检查页面是否包含表格
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        # 为表格创建新的幻灯片
                        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                        # 添加表格到幻灯片
                        rows, cols = len(table), len(table[0])
                        top = Inches(2)
                        left = Inches(1)
                        width = Inches(8)
                        height = Inches(0.8)
                        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
                        # 设置单元格文本
                        for r_idx, row in enumerate(table):
                            for c_idx, cell in enumerate(row):
                                table_shape.cell(r_idx, c_idx).text = cell               
                # 检查页面是否包含图片
                images = page.images
                if images:
                    for image in images:
                        # 为图片创建新的幻灯片
                        slide = ppt.slides.add_slide(ppt.slide_layouts[6])
                        im = Image.open(BytesIO(page.extract_image(image)['image']))
                        image_stream = BytesIO()
                        im.save(image_stream, format='PNG')
                        image_stream.seek(0)
                        left = Inches(1)
                        top = Inches(1)
                        slide.shapes.add_picture(image_stream, left, top, width=Inches(5.5), height=Inches(4.5))     
        # 创建临时文件保存pptx
        pptx_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx', mode='w+b')
        ppt.save(pptx_temp.name)
        
        # 返回pptx文件
        pptx_temp.seek(0)
        return send_file(pptx_temp.name, as_attachment=True, download_name='converted.pptx')
    finally:
        # 删除临时的PDF和PPTX文件
        os.unlink(pdf_temp.name)
        if pptx_temp:  # 检查pptx_temp是否已初始化
            os.unlink(pptx_temp.name)


if __name__ == '__main__':
    app.run(debug=False, host='0.0.0.0', port=5000)
